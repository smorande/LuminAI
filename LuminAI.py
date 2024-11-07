# ==========================================================================
# LuminaAI - Cognitive Interface
# ==========================================================================
# Developed by Dr. Swapnil M.
# Version: 1.0.0
# Creation Date: 2024
# 
# Open Source License:
# This program is free software: you can redistribute it and/or modify
# it under the terms of the MIT License as published by
# the Open Source Initiative.
#
# This program is distributed in the hope that it will be useful,
# but WITHOUT ANY WARRANTY; without even the implied warranty of
# MERCHANTABILITY or FITNESS FOR A PARTICULAR PURPOSE. See the
# MIT License for more details.
#
# Project Repository: https://github.com/yourusername/LuminaAI
# ==========================================================================

import streamlit as st
from openai import OpenAI
import os
import io
import base64
from PIL import Image
import pytesseract
from langchain_community.document_loaders import PyPDFLoader, TextLoader
from langchain_community.document_loaders import Docx2txtLoader as DocxLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Chroma
from textblob import TextBlob
import plotly.graph_objects as go
from sklearn.feature_extraction.text import TfidfVectorizer
import requests
from docx import Document
from pptx import Presentation
import csv
import json
import nltk
from nltk.translate.bleu_score import sentence_bleu, SmoothingFunction
import numpy as np
from dotenv import load_dotenv
from typing import List, Dict, Any, Generator, Optional, Tuple

# Load environment variables
load_dotenv()

class VectorStoreManager:
    """Manages document vectorization and similarity search with fallback options."""
    
    def __init__(self):
        self.vector_store = None
        self.vectorizer = None
        self.texts = None
        self.embeddings_matrix = None
        
    def create_vector_store(self, texts: List[Document], embeddings: np.ndarray) -> None:
        """
        Creates a vector store from documents using available backends
        """
        try:
            # Try using Chroma
            def embed_func(text):
                return self.vectorizer.transform([text]).toarray()[0]
            
            self.vector_store = Chroma.from_documents(
                documents=texts,
                embedding=embed_func
            )
        except ImportError:
            # Simple numpy-based fallback if Chroma is not available
            self._setup_numpy_fallback(texts, embeddings)

    def _setup_numpy_fallback(self, texts: List[Document], embeddings: np.ndarray) -> None:
        """
        Sets up a simple numpy-based vector store as last resort
        """
        self.texts = texts
        self.embeddings_matrix = embeddings
        
        class SimpleVectorStore:
            def __init__(self, parent):
                self.parent = parent
                
            def similarity_search(self, query: str, k: int = 3) -> List[Document]:
                query_vector = self.parent.vectorizer.transform([query]).toarray()[0]
                similarities = np.dot(self.parent.embeddings_matrix, query_vector)
                most_similar_indices = np.argsort(similarities)[-k:][::-1]
                return [self.parent.texts[i] for i in most_similar_indices]
        
        self.vector_store = SimpleVectorStore(self)

class PerformanceMetrics:
    """Handles calculation of various performance metrics for the AI responses"""
    
    def __init__(self):
        self.smoothing = SmoothingFunction()
    
    def calculate_perplexity(self, text: str) -> float:
        """
        Calculate perplexity score for given text
        """
        try:
            n = 3
            chars = list(text.lower())
            if len(chars) < n:
                return 0
                
            ngrams = [''.join(chars[i:i+n]) for i in range(len(chars)-n+1)]
            unique_ngrams = set(ngrams)
            ngram_counts = {ngram: ngrams.count(ngram) for ngram in unique_ngrams}
            
            log_prob = sum(np.log(ngram_counts[ngram]/len(ngrams)) for ngram in ngrams)
            perplexity = np.exp(-log_prob/len(ngrams))
            return min(perplexity, 1000)
        except Exception as e:
            print(f"Error calculating perplexity: {str(e)}")
            return 0
    
    def simple_tokenize(self, text: str) -> List[str]:
        """
        Simple tokenization without relying on NLTK punkt
        """
        # Remove special characters and extra whitespace
        cleaned = ''.join(c.lower() if c.isalnum() else ' ' for c in text)
        # Split on whitespace and filter out empty strings
        return [token for token in cleaned.split() if token]
    
    def calculate_bleu(self, reference: str, candidate: str) -> float:
        """
        Calculate BLEU score between reference and candidate texts
        using simple tokenization
        """
        try:
            if not reference or not candidate:
                return 0
            
            reference_tokens = self.simple_tokenize(reference)
            candidate_tokens = self.simple_tokenize(candidate)
            
            if not reference_tokens or not candidate_tokens:
                return 0
            
            # Calculate BLEU score with smoothing
            return sentence_bleu(
                [reference_tokens], 
                candidate_tokens,
                smoothing_function=self.smoothing.method1
            )
        except Exception as e:
            print(f"Error calculating BLEU score: {str(e)}")
            return 0

# Initialize OpenAI clients
nvidia_client = OpenAI(
    base_url=os.getenv('NVIDIA_API_BASE_URL'),
    api_key=os.getenv('NVIDIA_API_KEY')
)

openai_client = OpenAI(
    api_key=os.getenv('OPENAI_API_KEY')
)

# Initialize Claude client
claude_api_key = os.getenv('CLAUDE_API_KEY')

class LuminaAI:
    """Main AI engine class handling conversation and document processing"""
    
    def __init__(self):
        self.name = "Lumina"
        self.model = os.getenv('NVIDIA_MODEL')
        self.conversation_history: List[Dict[str, str]] = []
        self.document_store = None
        self.claude_api_key = os.getenv('CLAUDE_API_KEY')
        self.metrics = PerformanceMetrics()
        self.current_metrics = {
            'perplexity': 0,
            'bleu': 0
        }

    def update_metrics(self, response: str, reference: Optional[str] = None) -> None:
        """
        Update performance metrics based on the latest response
        """
        try:
            self.current_metrics['perplexity'] = self.metrics.calculate_perplexity(response)
            
            if reference:
                self.current_metrics['bleu'] = self.metrics.calculate_bleu(reference, response)
        except Exception as e:
            print(f"Error updating metrics: {str(e)}")
            self.current_metrics = {
                'perplexity': 0,
                'bleu': 0
            }

    def generate_image(self, prompt: str) -> Optional[str]:
        """
        Generate an image using Claude's API based on the prompt
        Returns the URL of the generated image or None if generation fails
        """
        url = "https://api.anthropic.com/v1/images"
        headers = {
            "Content-Type": "application/json",
            "X-API-Key": self.claude_api_key
        }
        data = {
            "model": os.getenv('CLAUDE_MODEL'),
            "prompt": prompt,
            "num_images": 1,
            "size": "1024x1024"
        }
        try:
            response = requests.post(url, headers=headers, json=data)
            response.raise_for_status()
            return response.json()['images'][0]['url']
        except Exception as e:
            st.error(f"Error generating image: {str(e)}")
            return None
        
    def generate_response(self, user_input: str) -> Generator[str, None, None]:
        """
        Generate AI response based on user input and conversation history
        Returns a generator for streaming responses
        """
        try:
            previous_response = self.conversation_history[-1]['content'] if self.conversation_history else ""
            
            messages = [
                *self.conversation_history,
                {"role": "user", "content": user_input}
            ]

            if self.document_store:
                relevant_docs = self.document_store.similarity_search(user_input, k=3)
                context = "\n".join([doc.page_content for doc in relevant_docs])
                messages.append({"role": "system", "content": f"Relevant context:\n{context}"})

            completion = nvidia_client.chat.completions.create(
                model=self.model,
                messages=messages,
                temperature=0.7,
                max_tokens=2048,
                stream=True
            )

            full_response = ""
            for chunk in completion:
                if chunk.choices[0].delta.content is not None:
                    full_response += chunk.choices[0].delta.content
                    yield chunk.choices[0].delta.content

            self.conversation_history.append({"role": "user", "content": user_input})
            self.conversation_history.append({"role": "assistant", "content": full_response})
            
            # Update metrics after response generation
            self.update_metrics(full_response, previous_response)
            
        except Exception as e:
            error_message = f"Error generating response: {str(e)}"
            st.error(error_message)
            yield error_message

    def set_document_store(self, document_store: Any) -> None:
        """Set the document store for context-aware responses"""
        self.document_store = document_store

    def clear_conversation(self) -> None:
        """Reset conversation history and metrics"""
        self.conversation_history = []
        self.current_metrics = {
            'perplexity': 0,
            'bleu': 0
        }

def detect_download_request(prompt: str) -> Dict[str, bool]:
    """
    Analyzes if user wants to download content in specific formats
    Returns a dictionary with boolean values for each supported format
    """
    default_formats = {
        "Image": False,
        "MS-Word": False,
        "CSV": False,
        "MS-PowerPoint": False,
        "Text": False,
        "Python": False
    }
    
    try:
        messages = [
            {
                "role": "system", 
                "content": """Analyze if the user wants to download or save content in specific formats.
                Return a JSON object with boolean values for these formats:
                Image, MS-Word, CSV, MS-PowerPoint, Text, Python."""
            },
            {"role": "user", "content": prompt}
        ]
        
        response = openai_client.chat.completions.create(
            model="gpt-4-0125-preview",
            messages=messages,
            temperature=0,
            max_tokens=150,
            response_format={"type": "json_object"}
        )
        
        try:
            result = json.loads(response.choices[0].message.content)
            if all(key in result for key in default_formats):
                return result
            return default_formats
            
        except json.JSONDecodeError:
            st.warning("Error parsing download options. Using defaults.")
            return default_formats
            
    except Exception as e:
        st.error(f"Error detecting download options: {str(e)}")
        return default_formats

def process_document(file) -> List[Document]:
    """
    Process uploaded documents and convert them to Document objects
    Supports PDF, DOCX, TXT, PY, and image files
    """
    file_extension = os.path.splitext(file.name)[1].lower()
    temp_file_path = f"temp{file_extension}"
    
    try:
        with open(temp_file_path, "wb") as temp_file:
            temp_file.write(file.getvalue())
        
        if file_extension == '.pdf':
            loader = PyPDFLoader(temp_file_path)
        elif file_extension == '.docx':
            loader = DocxLoader(temp_file_path)
        elif file_extension in ['.txt', '.py']:
            loader = TextLoader(temp_file_path)
        elif file_extension in ['.png', '.jpg', '.jpeg']:
            image = Image.open(temp_file_path)
            text = pytesseract.image_to_string(image)
            return [Document(page_content=text, metadata={"source": file.name})]
        else:
            raise ValueError(f"Unsupported file type: {file_extension}")
        
        documents = loader.load()
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        return text_splitter.split_documents(documents)
        
    except Exception as e:
        st.error(f"Error processing document {file.name}: {str(e)}")
        return []
        
    finally:
        if os.path.exists(temp_file_path):
            os.remove(temp_file_path)

def analyze_sentiment(text: str) -> float:
    """
    Analyze sentiment of text using TextBlob
    Returns a float between -1 (negative) and 1 (positive)
    """
    try:
        blob = TextBlob(text)
        return blob.sentiment.polarity
    except Exception as e:
        print(f"Error analyzing sentiment: {str(e)}")
        return 0.0

def update_document_store() -> None:
    """
    Update the document store with processed documents
    Creates and configures vector store for document similarity search
    """
    if not st.session_state.get('processed_documents'):
        return
        
    try:
        all_texts = [text for doc in st.session_state.processed_documents for text in doc['texts']]
        
        vectorizer = TfidfVectorizer()
        tfidf_matrix = vectorizer.fit_transform([text.page_content for text in all_texts])
        embeddings = tfidf_matrix.toarray()
        
        vector_store_manager = VectorStoreManager()
        vector_store_manager.vectorizer = vectorizer
        
        vector_store_manager.create_vector_store(all_texts, embeddings)
        st.session_state.lumina.set_document_store(vector_store_manager.vector_store)
        st.sidebar.success("Document store updated successfully.")
        
    except Exception as e:
        st.sidebar.error(f"Error updating document store: {str(e)}")
        st.sidebar.info("Try installing required packages: pip install chromadb")
        
def download_as_docx(content: str) -> bytes:
    """Convert content to DOCX format"""
    try:
        doc = Document()
        doc.add_paragraph(content)
        bio = io.BytesIO()
        doc.save(bio)
        return bio.getvalue()
    except Exception as e:
        st.error(f"Error creating Word document: {str(e)}")
        return b""

def download_as_pptx(content: str) -> bytes:
    """Convert content to PPTX format"""
    try:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide.shapes.add_textbox(10, 10, prs.slide_width-20, prs.slide_height-20)
        tf = txBox.text_frame
        tf.text = content
        bio = io.BytesIO()
        prs.save(bio)
        return bio.getvalue()
    except Exception as e:
        st.error(f"Error creating PowerPoint: {str(e)}")
        return b""

def download_as_csv(content: str) -> str:
    """Convert content to CSV format"""
    try:
        bio = io.StringIO()
        writer = csv.writer(bio)
        writer.writerow([content])
        return bio.getvalue()
    except Exception as e:
        st.error(f"Error creating CSV: {str(e)}")
        return ""

def needs_image_generation(prompt: str) -> bool:
    """Determine if prompt requires image generation"""
    try:
        image_keywords = ['draw', 'create', 'generate', 'make', 'produce', 'design']
        image_subjects = ['image', 'picture', 'diagram', 'figure', 'schematic', 'illustration', 'visual']
        prompt_lower = prompt.lower()
        return any(keyword in prompt_lower for keyword in image_keywords) and \
               any(subject in prompt_lower for subject in image_subjects)
    except Exception:
        return False

def display_performance_metrics(lumina: LuminaAI) -> None:
    """Display performance metrics in sidebar"""
    try:
        st.sidebar.subheader("🎯 Performance Metrics")
        
        perplexity = lumina.current_metrics['perplexity']
        st.sidebar.metric(
            "Perplexity",
            f"{perplexity:.2f}",
            help="Lower perplexity indicates more confident and fluent responses"
        )
        
        bleu = lumina.current_metrics['bleu']
        st.sidebar.metric(
            "BLEU Score",
            f"{bleu:.2f}",
            help="Higher BLEU score indicates better translation quality"
        )
    except Exception as e:
        st.sidebar.error(f"Error displaying metrics: {str(e)}")

# Streamlit UI Setup
st.set_page_config(page_title="LuminaAI", page_icon="🌟", layout="wide")

# Apply custom CSS styling
st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=Roboto:wght@300;400;700&display=swap');
    
    body {
        color: #E0E0E0;
        background-color: #1E1E1E;
        font-family: 'Roboto', sans-serif;
    }
    .stButton>button {
        color: #00BFFF;
        border-color: #00BFFF;
        border-radius: 20px;
        font-family: 'Roboto', sans-serif;
    }
    .stTextInput>div>div>input {
        color: #E0E0E0;
        background-color: #2E2E2E;
        border-radius: 20px;
        font-family: 'Roboto', sans-serif;
    }
    .sentiment-highlight {
        background-color: rgba(255, 255, 0, 0.3);
        padding: 2px 5px;
        border-radius: 3px;
    }
    .footer-container {
        display: flex;
        justify-content: space-between;
        align-items: center;
        margin-top: 2rem;
        padding: 1rem;
        border-top: 1px solid #333;
    }
    .app-header {
        text-align: center;
        margin-bottom: 2rem;
    }
    .chat-container {
        max-width: 800px;
        margin: 0 auto;
    }
</style>
""", unsafe_allow_html=True)

# Initialize session state
if 'lumina' not in st.session_state:
    st.session_state.lumina = LuminaAI()
if 'messages' not in st.session_state:
    st.session_state.messages = []
if 'download_options' not in st.session_state:
    st.session_state.download_options = {}
if 'processed_documents' not in st.session_state:
    st.session_state.processed_documents = []

lumina = st.session_state.lumina

# Main UI Layout
st.markdown('<div class="app-header">', unsafe_allow_html=True)
st.title("🌟 LuminaAI - Cognitive Interface")
st.markdown('</div>', unsafe_allow_html=True)

# Sidebar Configuration
with st.sidebar:
    st.title("LuminaAI Control Panel")
    
    st.subheader("🧠 LuminaAI Specs")
    st.markdown(f"**Model:** {lumina.model}")
    st.markdown("**Context Window:** 300,000 tokens")
    
    display_performance_metrics(lumina)
    
    st.subheader("📁 Document Upload")
    uploaded_files = st.file_uploader(
        "Upload documents for context (max 25)", 
        type=['pdf', 'docx', 'txt', 'py', 'png', 'jpg', 'jpeg'], 
        accept_multiple_files=True
    )

    if uploaded_files:
        if len(uploaded_files) > 25:
            st.warning("Maximum 25 files allowed. Only the first 25 will be processed.")
            uploaded_files = uploaded_files[:25]
        
        new_docs = [file for file in uploaded_files 
                   if file not in [doc['file'] for doc in st.session_state.processed_documents]]
        
        if new_docs:
            with st.spinner("Processing new documents..."):
                for file in new_docs:
                    try:
                        texts = process_document(file)
                        if texts:
                            st.session_state.processed_documents.append({'file': file, 'texts': texts})
                    except Exception as e:
                        st.error(f"Error processing {file.name}: {str(e)}")
            
            if st.session_state.processed_documents:
                st.success(f"{len(new_docs)} new document(s) processed successfully.")
                update_document_store()

    # Display processed documents
    if st.session_state.processed_documents:
        st.subheader("📚 Processed Documents")
        for doc in st.session_state.processed_documents:
            st.markdown(f"- {doc['file'].name}")

    # Conversation Analytics
    st.subheader("📊 Conversation Analytics")
    if st.session_state.messages:
        sentiments = [analyze_sentiment(msg["content"]) for msg in st.session_state.messages]
        
        fig_trend = go.Figure(data=go.Scatter(
            x=list(range(len(sentiments))),
            y=sentiments,
            mode='lines+markers',
            marker=dict(
                size=8,
                color=sentiments,
                colorscale='RdYlGn',
                showscale=True
            )
        ))
        fig_trend.update_layout(
            title='Sentiment Trend',
            xaxis_title='Message',
            yaxis_title='Sentiment',
            height=300,
            template="plotly_dark"
        )
        st.plotly_chart(fig_trend, use_container_width=True)
    else:
        st.info("Start a conversation to see analytics.")

# Main Chat Area
st.markdown('<div class="chat-container">', unsafe_allow_html=True)
st.subheader("💬 Interaction Panel")

# Display conversation history
for message in st.session_state.messages:
    with st.chat_message(message["role"]):
        st.markdown(message["content"])
        if "image_url" in message:
            st.image(message["image_url"], caption="Generated Image")

# Chat input and processing
if prompt := st.chat_input("What would you like to know or create?"):
    st.session_state.messages.append({"role": "user", "content": prompt})
    st.chat_message("user").markdown(prompt)

    with st.chat_message("assistant"):
        message_placeholder = st.empty()
        full_response = ""
        
        if needs_image_generation(prompt):
            with st.spinner("Generating image..."):
                image_url = lumina.generate_image(prompt)
            if image_url:
                st.image(image_url, caption="Generated Image")
                full_response = "I've generated an image based on your request. Here it is!"
                st.session_state.messages.append({
                    "role": "assistant", 
                    "content": full_response, 
                    "image_url": image_url
                })
            else:
                full_response = "I apologize, but I couldn't generate the image. Please try a different prompt."
                st.error(full_response)
        else:
            try:
                for chunk in lumina.generate_response(prompt):
                    full_response += chunk
                    message_placeholder.markdown(full_response + "▌")
                message_placeholder.markdown(full_response)
            except Exception as e:
                error_message = f"Error generating response: {str(e)}"
                st.error(error_message)
                full_response = error_message
        
        if not needs_image_generation(prompt):
            st.session_state.messages.append({
                "role": "assistant", 
                "content": full_response
            })
        
        # Display sentiment analysis
        sentiment = analyze_sentiment(full_response)
        st.markdown(
            f'<span class="sentiment-highlight">Response sentiment: {sentiment:.2f}</span>', 
            unsafe_allow_html=True
        )

        # Handle download options
        download_options = detect_download_request(prompt)
        if any(download_options.values()):
            st.subheader("Download Options")
            cols = st.columns(6)
            
            if download_options.get("Text", False):
                with cols[0]:
                    st.download_button(
                        "Download Text",
                        full_response,
                        file_name="response.txt",
                        mime="text/plain"
                    )
            
            if download_options.get("MS-Word", False):
                with cols[1]:
                    st.download_button(
                        "Download Word",
                        download_as_docx(full_response),
                        file_name="response.docx",
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                    )
            
            if download_options.get("CSV", False):
                with cols[2]:
                    st.download_button(
                        "Download CSV",
                        download_as_csv(full_response),
                        file_name="response.csv",
                        mime="text/csv"
                    )
            
            if download_options.get("MS-PowerPoint", False):
                with cols[3]:
                    st.download_button(
                        "Download PowerPoint",
                        download_as_pptx(full_response),
                        file_name="response.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
            
            if download_options.get("Python", False):
                with cols[4]:
                    st.download_button(
                        "Download Python",
                        full_response,
                        file_name="response.py",
                        mime="text/x-python"
                    )

            if download_options.get("Image", False) and "image_url" in st.session_state.messages[-1]:
                with cols[5]:
                    response = requests.get(st.session_state.messages[-1]["image_url"])
                    if response.status_code == 200:
                        st.download_button(
                            "Download Image",
                            response.content,
                            file_name="generated_image.png",
                            mime="image/png"
                        )
st.markdown('</div>', unsafe_allow_html=True)

# Footer
st.markdown('''
<div class="footer-container">
    <div>Powered by LuminaAI - Illuminating Your Path to Knowledge</div>
    <div>
''', unsafe_allow_html=True)

if st.button(
    "Clear Screen", 
    key="clear_screen", 
    help="Reset the conversation context", 
    type="secondary",
    use_container_width=True
):
    st.session_state.messages = []
    st.session_state.download_options = {}
    lumina.clear_conversation()
    st.rerun()

st.markdown('</div></div>', unsafe_allow_html=True)